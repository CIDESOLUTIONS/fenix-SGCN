from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO

# Helper functions
def add_title(prs, text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    title = slide.shapes.title
    if not title:
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    else:
        title.left, title.top, title.width, title.height = Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    return slide

def add_bullets(slide, items, left=0.5, top=1.2, width=9, height=5):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for i, it in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = it
        p.level = 0
        p.font.size = Pt(20)
    return box

def add_chevron_row(slide, labels, top_in=2.0, chevron_w=2.0, gap=0.2):
    left = 0.5
    for i, label in enumerate(labels):
        shape = slide.shapes.add_shape(MSO.CHEVRON, Inches(left), Inches(top_in), Inches(chevron_w), Inches(1.2))
        shape.text_frame.text = label
        shape.fill.solid()
        # alternate neutral colors
        if i % 2 == 0:
            shape.fill.fore_color.rgb = RGBColor(224, 224, 224)
        else:
            shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
        shape.line.color.rgb = RGBColor(120, 120, 120)
        left += chevron_w - 0.3 + gap  # slight overlap for chevron effect

def add_kpi_cards(slide, cards):
    left = 0.5
    top = 1.5
    card_w = 3.0
    card_h = 1.5
    for title, desc in cards:
        rect = slide.shapes.add_shape(MSO.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(card_w), Inches(card_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(245, 245, 245)
        rect.line.color.rgb = RGBColor(180, 180, 180)
        tf = rect.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.LEFT
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.level = 1
        left += card_w + 0.3
        if left + card_w > 10:
            left = 0.5
            top += card_h + 0.3

def add_two_column(slide, left_title, left_points, right_title, right_points):
    # Left column
    ltitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4.5), Inches(0.5))
    ltitle.text_frame.text = left_title
    ltitle.text_frame.paragraphs[0].font.bold = True
    ltitle.text_frame.paragraphs[0].font.size = Pt(20)
    add_bullets(slide, left_points, left=0.6, top=1.5, width=4.3, height=4)
    # Right column
    rtitle = slide.shapes.add_textbox(Inches(5.2), Inches(1.0), Inches(4.5), Inches(0.5))
    rtitle.text_frame.text = right_title
    rtitle.text_frame.paragraphs[0].font.bold = True
    rtitle.text_frame.paragraphs[0].font.size = Pt(20)
    add_bullets(slide, right_points, left=5.3, top=1.5, width=4.3, height=4)

# Build presentation
prs = Presentation()

# 1. Cover
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Sensibilización en Continuidad del Negocio"
slide.placeholders[1].text = "ISO 22301 – Sistema de Gestión de Continuidad del Negocio (SGCN)\nCliente financiero de gobierno – Colombia"

# 2. Qué es / Qué no es
slide = add_title(prs, "¿Qué es y qué NO es un SGCN (ISO 22301)?")
add_two_column(
    slide,
    "Es",
    [
        "Un sistema de gestión (BCMS) basado en riesgos y requisitos ISO 22301.",
        "Marco para planear, operar, evaluar y mejorar continuidad (PDCA).",
        "Incluye BIA, evaluación de riesgos, estrategias, planes, pruebas y mejora.",
        "Enfocado en proteger servicios y objetivos misionales."
    ],
    "No es",
    [
        "No es solo TI ni solo recuperación de desastres.",
        "No es un documento estático ni solo plan de emergencias.",
        "No es exclusivo del área de riesgos: es transversal.",
        "No es costoso por sí mismo: es inversión en resiliencia."
    ]
)

# 3. Flujo de análisis (alto nivel)
slide = add_title(prs, "Flujo de Análisis del SGCN")
add_chevron_row(
    slide,
    [
        "Seleccionar procesos críticos",
        "ARA – Análisis de Riesgos",
        "BIA – Impacto al Negocio",
        "Estrategias de continuidad",
        "Planes (BCP/DRP/CRP)",
        "Plan de pruebas",
        "Correcciones y cambios",
        "Ciclo PDCA"
    ],
    top_in=1.6,
    chevron_w=2.1
)

# 4. BIA – resultados clave
slide = add_title(prs, "BIA: resultados y decisiones clave")
add_two_column(
    slide,
    "Parámetros del BIA",
    [
        "Servicios/actividades críticas y dependencias.",
        "MAO/MTPD (Máx. tiempo tolerable de interrupción).",
        "RTO (objetivo de tiempo de recuperación).",
        "RPO (objetivo de punto de recuperación).",
        "Impactos (financieros, legales, reputacionales, operativos, ciudadanos)."
    ],
    "Decisiones que habilita",
    [
        "Priorización de procesos y recursos.",
        "Niveles de continuidad requeridos (tiers).",
        "Criterios de activación y escalamiento.",
        "Bases para estrategias y dimensionamiento de soluciones."
    ]
)

# 5. ARA – evaluación de riesgos
slide = add_title(prs, "ARA: Evaluación de Riesgos para Continuidad")
add_two_column(
    slide,
    "Alcance y entradas",
    [
        "Contexto (interno/externo) y partes interesadas.",
        "Catálogo de amenazas y vulnerabilidades.",
        "Controles existentes y riesgos residuales.",
        "Criterios de aceptabilidad (apetito/tolerancia)."
    ],
    "Resultados",
    [
        "Mapa de riesgos por proceso/servicio.",
        "Escenarios de interrupción priorizados.",
        "Tratamientos de riesgo y controles preventivos.",
        "Insumos para estrategias de continuidad."
    ]
)

# 6. Estrategias de continuidad
slide = add_title(prs, "Selección de Estrategias de Continuidad")
add_bullets(
    slide,
    [
        "Personas: reemplazos, multifuncionalidad, turnos, contacto.",
        "Sitios: trabajo remoto, sitio alterno, acuerdos recíprocos.",
        "Tecnología: alta disponibilidad, replicación, DR en nube, backups.",
        "Proveedores: contratos, SLA, multisourcing.",
        "Datos: clasificación, RPO/RTO, retención y restauración.",
        "Procesos: manuales alternos, colas, degradación controlada."
    ]
)

# 7. Planes de respuesta y gestión de crisis
slide = add_title(prs, "Planes y Gestión de Crisis")
add_two_column(
    slide,
    "Tipos de planes",
    [
        "BCP por proceso/servicio crítico.",
        "DRP (recuperación tecnológica).",
        "CMP/CRP (gestión de crisis y comunicaciones).",
        "Plan de continuidad de proveedores."
    ],
    "Activación y coordinación",
    [
        "Criterios de activación claros y roles definidos.",
        "Centro de crisis / Comité de continuidad.",
        "Comunicación interna/externa (stakeholders, entes de control).",
        "Retorno a la operación normal y lecciones aprendidas."
    ]
)

# 8. Programa de pruebas y ejercicios
slide = add_title(prs, "Programa de Pruebas y Ejercicios")
add_bullets(
    slide,
    [
        "Modalidades: revisión de escritorio, simulacro, técnico, corte de servicios.",
        "Cobertura mínima anual de procesos críticos y tecnología habilitante.",
        "Criterios de éxito: cumplimiento de RTO/RPO, coordinación, comunicaciones.",
        "Registro de hallazgos, plan de acción y gestión de cambios."
    ]
)

# 9. Implementación por fases (hoja de ruta)
slide = add_title(prs, "Cómo se implementa (Hoja de Ruta)")
add_bullets(
    slide,
    [
        "1) Gobierno y alcance: política, objetivos, alcance, roles, recursos.",
        "2) Contexto e interesados (Cláus. 4): requisitos internos/externos.",
        "3) BIA y ARA (Cláus. 8.2): procesos críticos, RTO/RPO, riesgos.",
        "4) Estrategias y soluciones (Cláus. 8.3).",
        "5) Planes y procedimientos (Cláus. 8.4).",
        "6) Soporte (Cláus. 7): competencias, conciencia, comunicación.",
        "7) Operación (Cláus. 8) + Preparación y respuesta a incidentes.",
        "8) Evaluación del desempeño (Cláus. 9): KPIs, auditoría interna.",
        "9) Mejora (Cláus. 10): no conformidades, acciones correctivas.",
        "10) Preparación para certificación (opcional)."
    ]
)

# 10. Cómo se mantiene (operación continua)
slide = add_title(prs, "Cómo se mantiene (Operación Continua)")
add_kpi_cards(
    slide,
    [
        ("Gobierno", "Revisión por la dirección semestral; política y objetivos vigentes."),
        ("Cambios", "Gestión de cambios al BCMS (tecnología, procesos, proveedores)."),
        ("Entrenamiento", "% personal crítico capacitado; inducción y refrescos."),
        ("Ejercicios", "% procesos críticos probados; cumplimiento RTO/RPO."),
        ("Auditoría", "Hallazgos cerrados dentro de plazo; eficacia de acciones."),
        ("Indicadores", "Disponibilidad de servicios, tiempo de activación de crisis.")
    ]
)

# 11. ISO 22301 – estructura y evolución
slide = add_title(prs, "ISO 22301: estructura y evolución")
add_two_column(
    slide,
    "Estructura (Anexo SL)",
    [
        "Cláusulas 4–10: Contexto, Liderazgo, Planificación, Soporte, Operación, Evaluación, Mejora.",
        "Mayor énfasis en desempeño y gestión del cambio (versión 2019).",
        "Documentación basada en 'información documentada'."
    ],
    "Evolución 2012 → 2019",
    [
        "Alineación plena con Anexo SL.",
        "Racionalización de requisitos y enfoque en resultados.",
        "Centralización de requisitos operativos en Cláusula 8."
    ]
)

# 12. Roles y responsabilidades
slide = add_title(prs, "Roles y Responsabilidades")
add_bullets(
    slide,
    [
        "Alta Dirección: liderazgo, política, recursos y apetito de riesgo.",
        "Propietario del BCMS / Oficial de Continuidad: coordina el sistema.",
        "Comité de Continuidad/Crisis: decisiones y escalamiento.",
        "Dueños de Proceso: BIA, planes y prueba de su proceso.",
        "TI/Seguridad/Ciber: DR, alta disponibilidad, ciber-respuesta.",
        "Comunicaciones: relación con entes de control y ciudadanía.",
        "Terceros/Proveedores: cumplimiento de SLA y continuidad."
    ]
)

# 13. Beneficios para entidad financiera de gobierno
slide = add_title(prs, "Beneficios de la adopción (Entidad financiera de gobierno)")
add_bullets(
    slide,
    [
        "Resiliencia de servicios misionales y continuidad a la ciudadanía.",
        "Cumplimiento normativo (SFC, MIPG, seguridad digital).",
        "Reducción de pérdidas financieras y reputacionales.",
        "Mejor gestión de terceros y ciberresiliencia.",
        "Mejora de confianza de órganos de control y calificadoras."
    ]
)

# 14. Alineación normativa Colombia (referencial)
slide = add_title(prs, "Alineación normativa – Colombia (referencial)")
add_bullets(
    slide,
    [
        "SFC – SARO (CE 041/2007) y lineamientos sobre continuidad y ciberseguridad.",
        "MIPG – Gestión del riesgo, continuidad y seguridad digital en entidades públicas.",
        "Buenas prácticas complementarias: ISO 31000 (riesgos) e ISO/IEC 27001 (seguridad).",
        "Armonización con planes sectoriales (MinTIC, DAFP) cuando aplique."
    ]
)

# 15. Próximos pasos
slide = add_title(prs, "Próximos pasos y plan de trabajo")
add_bullets(
    slide,
    [
        "Validar alcance, procesos críticos y patrocinios.",
        "Plan de entrevistas para BIA/ARA y levantamiento de dependencias.",
        "Definir criterios de continuidad (RTO/RPO/MAO) y niveles de servicio.",
        "Seleccionar estrategias y dimensionar soluciones (incluye CAPEX/OPEX).",
        "Diseñar e implantar planes; formar equipos y realizar ejercicios.",
        "Definir KPIs, auditoría interna y calendario de mantenimiento."
    ]
)

# Save
path = "/c/users/meciz/documents/fenix-SGCN/SGCN_ISO22301_Detallada.pptx"
prs.save(path)
path
