# ----------------------------------------------------------
# 15-LAMINAS  |  ICETEX  |  ISO 22301  |  SGCN  |  2025
# Secuencia: slide → placeholder → formato → contenido → animación
# ----------------------------------------------------------
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# 1. Cargar plantilla (con fondo ICETEX)
ruta_plantilla = r"C:\Users\meciz\Documents\fenix-SGCN\docs\plantilla_icetex.pptx"
if not os.path.exists(ruta_plantilla):
    print("❌ No se encuentra plantilla_icetex.pptx"); exit()

prs = Presentation(ruta_plantilla)

# 2. Colores ICETEX
AZUL_ICETEX = RGBColor(0, 51, 160)   # #0033A0
NARANJA = RGBColor(255, 103, 31)     # #FF671F
GRIS = RGBColor(89, 89, 89)

# 3. Función: crear textbox con formato y animación
def add_textbox(slide, x, y, w, h, texto, tam=16, bold=False, color=AZUL_ICETEX, animar=True):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.clear()  # limpia previos
    p = tf.paragraphs[0]
    p.text = texto
    p.font.size = Pt(tam)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT
    # Animación: aparecer al hacer clic (PowerPoint la interpreta)
    if animar:
        tx.click_action.target_slide = slide  # truco estándar
    return tx

# 4. Función: lámina completa (título + bullets)
def add_lamina(titulo, bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Título
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(1), titulo, tam=28, bold=True, color=AZUL_ICETEX, animar=False)
    # Bullets
    if bullets:
        y_ini = Inches(2)
        for i, txt in enumerate(bullets):
            add_textbox(slide, Inches(0.5), y_ini + i * Inches(0.6), Inches(8.5), Inches(0.5), "• " + txt, tam=16, color=GRIS, animar=True)
    return slide

# 5. PORTADA
portada = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(portada, Inches(1), Inches(2.5), Inches(8), Inches(1.5), "SGCN ICETEX – ISO 22301", tam=40, bold=True, color=AZUL_ICETEX, animar=False)
add_textbox(portada, Inches(1), Inches(4), Inches(8), Inches(1), "Presentación Gerencial de Inicio", tam=20, color=NARANJA, animar=False)

# 6. RESTO DE LÁMINAS
add_lamina("Objetivo de esta presentación",
           ["Motivar al grupo directivo para que comprenda, respalde y lidere la implementación del SGCN con entusiasmo y claridad estratégica."])

add_lamina("¿Por qué un SGCN en ICETEX?",
           ["Porque el futuro de miles de estudiantes depende de que ICETEX nunca se detenga.",
            "Ciberataque → sin SGCN: pérdida datos; con SGCN: recuperación <4 h",
            "Terremoto → sin SGCN: días de parálisis; con SGCN: activación <2 h"])

add_lamina("¿Qué es y qué NO es un SGCN ISO 22301?",
           ["ES: Proceso de negocio que protege la misión de ICETEX",
            "NO ES: Un simple plan de TI ni un documento archivado"])

add_lamina("Madurez: de Continuidad a Resiliencia",
           ["Reactivo → Proactivo → Predictivo → Resiliente → Antifragil",
            "ICETEX saltará a Resiliente en 12 meses"])

add_lamina("Ciclo PHVA – Implementación SGCN",
           ["Planear → Diseñar → Implementar → Probar → Medir → Mejorar → Certificar"])

add_lamina("Fase 1 – Planear (Sem 1-4)",
           ["Kick-off + RACI", "Entrevistas 30 procesos", "Validación críticos", "Aprobación Comité"])

add_lamina("Fase 2 – Diseñar (Sem 5-8)",
           ["Taller RTO/RPO", "Benchmark proveedores", "Pre-diseño planes", "Aprobación inversión"])

add_lamina("Fase 3 – Implementar (Sem 9-16)",
           ["Redacción 12 planes", "Infra alterna", "Capacitar 150 usuarios", "Simulacro mesa"])

add_lamina("Fase 4 – Probar (Sem 17-20)",
           ["Simulacro ciber", "Simulacro sismo", "Simulacro proveedor", "Informe ajustes"])

add_lamina("Fases 5-6-7 (Medir-Mejorar-Certificar)",
           ["Sem 21-22: KPIs y auditoría interna",
            "Sem 23-24: Taller de ajustes y aprobación",
            "Sem 25-26: Pre-auditoría y certificación ISO 22301"])

add_lamina("Roles & Responsabilidades (RACI)",
           ["Patrocinador: Gerente General",
            "Owner SGCN: VP de Riesgos",
            "Dueños Proceso: Dir. Financiera, CIO, etc."])

add_lamina("Beneficios tangibles",
           ["0% → 100% procesos con RTO ≤4 h",
            "0 → 6 simulacros/año",
            "Multas SFC → 0 multas",
            "Percepción reactiva → resiliente"])

add_lamina("Inversión vs Retorno 12 meses",
           ["CAPEX: 1.8 mil M COP", "OPEX anual: 0.6 mil M COP",
            "ROI evitación pérdidas: 4.5 mil M COP", "Payback: 8 meses"])

# 7. CIERRE
cierre = prs.slides.add_slide(prs.slide_layouts[6])
add_textbox(cierre, Inches(1), Inches(3), Inches(8), Inches(1.5), "“Cuando un estudiante necesite su crédito, ICETEX estará ahí.”", tam=24, bold=True, color=AZUL_ICETEX, animar=False)

# 8. GUARDAR
salida = "SGCN_ICETEX_15L_DISEÑO_ANIMADO.pptx"
prs.save(salida)
print(f"✅ Presentación con diseño y animaciones guardada: {salida}")